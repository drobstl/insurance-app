#!/usr/bin/env python3
"""AgentForLife — Premium ARCH Grants Pitch Deck Generator (Quility Colors)"""

import os, tempfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════════════════════════
#  DIMENSIONS
# ═══════════════════════════════════════════════════════════════
SW = Emu(12192000)   # 13.333 in
SH = Emu(6858000)    # 7.5 in
M  = Inches(0.8)     # margin
FH = Inches(0.45)    # footer height

# ═══════════════════════════════════════════════════════════════
#  QUILITY COLOR PALETTE
# ═══════════════════════════════════════════════════════════════
C_PRI      = RGBColor(0x00, 0x58, 0x51)
C_PRI_DK   = RGBColor(0x00, 0x46, 0x41)
C_PRI_DKST = RGBColor(0x0E, 0x36, 0x31)
C_BG1      = RGBColor(0x0E, 0x36, 0x31)
C_BG2      = RGBColor(0x05, 0x18, 0x16)
C_ACC      = RGBColor(0x46, 0xC3, 0xB2)
C_ACC_LT   = RGBColor(0x6E, 0xCF, 0x93)
C_BLUE     = RGBColor(0x00, 0x83, 0xEB)
C_YEL      = RGBColor(0xFF, 0xCC, 0x00)
C_CORAL    = RGBColor(0xFF, 0x58, 0x51)
C_W        = RGBColor(0xFF, 0xFF, 0xFF)
C_OW       = RGBColor(0xE0, 0xE0, 0xE0)
C_LT       = RGBColor(0xBB, 0xBB, 0xBB)
C_MD       = RGBColor(0x88, 0x88, 0x88)
C_DIM      = RGBColor(0x66, 0x66, 0x66)
C_CARD     = RGBColor(0x0A, 0x30, 0x2C)
C_CARD_B   = RGBColor(0x18, 0x4A, 0x44)
C_FT       = RGBColor(0x04, 0x12, 0x10)
C_RED_CARD = RGBColor(0x2E, 0x12, 0x12)
C_RED_BRD  = RGBColor(0x5A, 0x20, 0x20)
C_BEZEL    = RGBColor(0x14, 0x14, 0x14)
FONT       = "Montserrat"

# ═══════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════

def _grad(shape, c1, c2, ang=90):
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0
    gf = shape._element.find(".//" + qn("a:gradFill"))
    if gf is not None:
        for t in (qn("a:lin"), qn("a:path"), qn("a:tileRect")):
            for e in gf.findall(t):
                gf.remove(e)
        lin = etree.SubElement(gf, qn("a:lin"))
        lin.set("ang", str(ang * 60000))
        lin.set("scaled", "0")

def bg(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _grad(s, C_BG1, C_BG2, 150)
    s.line.fill.background()
    sp = s._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)

def ft(slide, text="agentforlife.app  |  Brainstorm Labs LLC  |  St. Louis, MO"):
    y = SH - FH
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SW, FH)
    bar.fill.solid(); bar.fill.fore_color.rgb = C_FT; bar.line.fill.background()
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SW, Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_ACC; ln.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(8)
    r = p.add_run(); r.text = text; r.font.size = Pt(10); r.font.color.rgb = C_MD; r.font.name = FONT

def tx(slide, text, x, y, w, h, sz=14, c=C_W, b=False, al=PP_ALIGN.LEFT, va=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = va
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b; r.font.name = FONT
    return tb

def mtx(slide, lines, x, y, w, h, sz=14, c=C_W, b=False, al=PP_ALIGN.LEFT, sp=Pt(6)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al; p.space_after = sp
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b; r.font.name = FONT
    return tb

def crd(slide, x, y, w, h, fill=C_CARD, brd=C_CARD_B):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = brd; s.line.width = Pt(1)
    try: s.adjustments[0] = 0.04
    except: pass
    return s

def acrd(slide, x, y, w, h, ac=C_ACC):
    c = crd(slide, x, y, w, h)
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(4))
    b.fill.solid(); b.fill.fore_color.rgb = ac; b.line.fill.background()
    return c

def lcrd(slide, x, y, w, h, ac=C_ACC):
    c = crd(slide, x, y, w, h)
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(5), h)
    b.fill.solid(); b.fill.fore_color.rgb = ac; b.line.fill.background()
    return c

def circ(slide, x, y, sz, fc=C_ACC, t="", tc=C_W, ts=14):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = fc; s.line.fill.background()
    if t:
        tf = s.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t
        r.font.size = Pt(ts); r.font.color.rgb = tc; r.font.bold = True; r.font.name = FONT
    return s

def ttl(slide, title, sub=None):
    tx(slide, title, M, Inches(0.4), Inches(10), Inches(0.6), sz=36, b=True)
    if sub:
        tx(slide, sub, M, Inches(1.0), Inches(11), Inches(0.4), sz=16, c=C_ACC)

def bullet(slide, text, x, y, w, bc=C_ACC, tc=C_OW, sz=13):
    d = Inches(0.12)
    circ(slide, x, y + Inches(0.04), d, bc)
    tx(slide, text, x + Inches(0.25), y, w - Inches(0.25), Inches(0.35), sz=sz, c=tc)

def stat(slide, val, label, x, y, w=Inches(3.5), vc=C_ACC, vs=48, src=None):
    tx(slide, val, x, y, w, Inches(0.8), sz=vs, c=vc, b=True, al=PP_ALIGN.CENTER)
    tx(slide, label, x, y + Inches(0.7), w, Inches(0.4), sz=12, c=C_OW, al=PP_ALIGN.CENTER)
    if src:
        tx(slide, src, x, y + Inches(1.1), w, Inches(0.25), sz=8, c=C_DIM, al=PP_ALIGN.CENTER)

def phone_frame(slide, img_path, x, y, w=Inches(2.0), h=Inches(4.1)):
    pad = Inches(0.07)
    crd(slide, x - pad, y - pad, w + pad * 2, h + pad * 2, C_BEZEL, C_ACC)
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y, w, h)

def dash_frame(slide, img_path, x, y, w=Inches(4.8), h=Inches(1.2)):
    pad = Inches(0.05)
    crd(slide, x - pad, y - pad, w + pad * 2, h + pad * 2, C_BEZEL, C_CARD_B)
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y, w, h)

# ═══════════════════════════════════════════════════════════════
#  IMAGE EXTRACTION
# ═══════════════════════════════════════════════════════════════

def extract_images(src, out):
    prs = Presentation(src)
    imgs = {}
    for si, slide in enumerate(prs.slides):
        ii = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:
                try:
                    blob = shape.image.blob
                    ext = shape.image.content_type.split("/")[-1].replace("jpeg", "jpg")
                    p = os.path.join(out, f"s{si+1}_i{ii}.{ext}")
                    with open(p, "wb") as f: f.write(blob)
                    imgs[(si + 1, ii)] = p
                    ii += 1
                except: ii += 1
    return imgs

# ═══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def s01_title(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)

    circ(sl, M, Inches(0.4), Inches(0.65), C_PRI, "\u221E", C_ACC, 30)
    tx(sl, "AgentForLife", Inches(1.65), Inches(0.47), Inches(4), Inches(0.55), sz=32, b=True)

    for i, (line, color) in enumerate([
        ("Kill Chargebacks.", C_W),
        ("Explode Referrals.", C_W),
        ("Triple Income.", C_YEL),
    ]):
        tx(sl, line, M, Inches(1.55) + Inches(i * 0.7), Inches(6), Inches(0.7), sz=42, c=color, b=True)

    tx(sl, "AI-powered retention, referrals & rewrites\nfor insurance agents.",
       M, Inches(3.75), Inches(5.5), Inches(0.8), sz=16, c=C_OW)

    badge = crd(sl, M, Inches(4.8), Inches(4.5), Inches(0.4), C_PRI_DK, C_ACC)
    tx(sl, "LIVE ON iOS & ANDROID  \u2022  app.agentforlife.com",
       M + Inches(0.15), Inches(4.8), Inches(4.2), Inches(0.4), sz=10, c=C_ACC, al=PP_ALIGN.LEFT, va=MSO_ANCHOR.MIDDLE)

    px1, px2 = Inches(7.8), Inches(10.2)
    py, pw, ph = Inches(0.6), Inches(2.0), Inches(4.2)
    phone_frame(sl, im.get((4, 1)), px1, py, pw, ph)
    phone_frame(sl, im.get((1, 1)), px2, py, pw, ph)
    tx(sl, "CLIENT APP", px1, py + ph + Inches(0.12), pw, Inches(0.25), sz=9, c=C_LT, al=PP_ALIGN.CENTER)
    tx(sl, "AI REFERRAL TEXT", px2, py + ph + Inches(0.12), pw, Inches(0.25), sz=9, c=C_LT, al=PP_ALIGN.CENTER)

    ft(sl)


def s02_problem(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "The Problem", "Agents buy a full tank of gas \u2014 but their engine leaks fuel.")

    cw = Inches(3.6)
    gap = Inches(0.2)
    x0 = M
    cy = Inches(1.6)
    ch = Inches(2.2)
    for i, (v, lbl, src) in enumerate([
        ("84%", "average agency retention\u2014 16% of book lost every year",
         "Source: InsuredMine, Nationwide"),
        ("7\u20139x", "more expensive to acquire\na new customer than retain one",
         "Source: Nationwide, Glassbox, bolttech"),
        ("83%", "of clients would refer\u2014\nbut only 29% actually do",
         "Source: Texas Tech University"),
    ]):
        cx = x0 + Inches(i * 3.9)
        acrd(sl, cx, cy, cw, ch, [C_ACC, C_BLUE, C_YEL][i])
        tx(sl, v, cx, cy + Inches(0.2), cw, Inches(0.7), sz=44, c=[C_ACC, C_BLUE, C_YEL][i], b=True, al=PP_ALIGN.CENTER)
        tx(sl, lbl, cx + Inches(0.3), cy + Inches(0.95), cw - Inches(0.6), Inches(0.6), sz=12, c=C_OW, al=PP_ALIGN.CENTER)
        tx(sl, src, cx + Inches(0.3), cy + Inches(1.7), cw - Inches(0.6), Inches(0.25), sz=8, c=C_DIM, al=PP_ALIGN.CENTER)

    line_y = Inches(4.1)
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, line_y, Inches(11.7), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_CARD_B; ln.line.fill.background()

    lcrd(sl, M, Inches(4.35), Inches(5.5), Inches(2.3), C_CORAL)
    tx(sl, "The Missing Revenue", M + Inches(0.25), Inches(4.45), Inches(5), Inches(0.35), sz=16, c=C_W, b=True)
    lines = [
        "An agent buys 20 leads/week \u2192 ceiling: 20 policies",
        "Real potential with referrals & rewrites: 60 policies",
        "That\u2019s $2,400/lead left on the table",
        "$200B in life insurance lapsed or surrendered annually",
    ]
    for j, l in enumerate(lines):
        bullet(sl, l, M + Inches(0.2), Inches(4.9) + Inches(j * 0.38), Inches(5.5), C_CORAL)

    lcrd(sl, Inches(7.0), Inches(4.35), Inches(5.5), Inches(2.3), C_YEL)
    tx(sl, "Why Agents Leave Money on the Table", Inches(7.25), Inches(4.45), Inches(5), Inches(0.35), sz=16, c=C_W, b=True)
    pains = [
        "No automated retention \u2014 clients forget their agent in weeks",
        "No referral system \u2014 agents rely on hope, not tools",
        "No rewrite alerts \u2014 anniversary renewals go unmonitored",
        "90% of new agents quit within 3 years (LIMRA)",
    ]
    for j, p in enumerate(pains):
        bullet(sl, p, Inches(7.25), Inches(4.9) + Inches(j * 0.38), Inches(5.2), C_YEL)

    ft(sl)


def s03_solution(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "The Solution", "Three revenue streams. One branded app. Powered by AI.")

    cols = [
        (C_ACC,  "R", "Retention", "7+", "auto touchpoints/yr",
         "Birthday, holiday, and anniversary messages.\nAI conservation catches lapsing policies within hours."),
        (C_BLUE, "F", "Referrals", "1-Tap", "referral system",
         "Client picks a contact, warm intro text goes out.\nAI business line qualifies and books the appointment."),
        (C_YEL,  "W", "Rewrites",  "1-Year", "anniversary outreach",
         "Clients notified before renewal with a booking link.\nRevenue agents never knew existed."),
    ]

    cw = Inches(3.6)
    cy = Inches(1.6)
    ch = Inches(4.7)

    for i, (ac, icon, name, big, sub, desc) in enumerate(cols):
        cx = M + Inches(i * 3.9)
        acrd(sl, cx, cy, cw, ch, ac)
        circ(sl, cx + Inches(0.3), cy + Inches(0.35), Inches(0.55), ac, icon, C_W, 20)
        tx(sl, name, cx + Inches(1.0), cy + Inches(0.4), Inches(2.2), Inches(0.45), sz=22, c=C_W, b=True)
        tx(sl, big, cx, cy + Inches(1.2), cw, Inches(0.7), sz=48, c=ac, b=True, al=PP_ALIGN.CENTER)
        tx(sl, sub, cx, cy + Inches(1.9), cw, Inches(0.35), sz=14, c=C_OW, al=PP_ALIGN.CENTER)
        ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx + Inches(0.4), cy + Inches(2.45), cw - Inches(0.8), Pt(1))
        ln.fill.solid(); ln.fill.fore_color.rgb = C_CARD_B; ln.line.fill.background()
        tx(sl, desc, cx + Inches(0.35), cy + Inches(2.7), cw - Inches(0.7), Inches(1.5), sz=12, c=C_LT)

    ft(sl)


def s04_product(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "The Product \u2014 Live Today")

    tx(sl, "CLIENT MOBILE APP", M, Inches(1.25), Inches(5), Inches(0.3), sz=12, c=C_ACC, b=True)
    pw, ph = Inches(1.65), Inches(3.35)
    py = Inches(1.65)
    phones = [
        (im.get((4, 0)), "Agent Profile"),
        (im.get((4, 1)), "Home Screen"),
        (im.get((4, 2)), "Policy View"),
    ]
    for i, (img, lbl) in enumerate(phones):
        px = M + Inches(i * 2.0)
        phone_frame(sl, img, px, py, pw, ph)
        tx(sl, lbl, px, py + ph + Inches(0.1), pw, Inches(0.25), sz=9, c=C_LT, al=PP_ALIGN.CENTER)

    dx = Inches(7.0)
    tx(sl, "AGENT WEB DASHBOARD", dx, Inches(1.25), Inches(5.5), Inches(0.3), sz=12, c=C_ACC, b=True)
    dw, dh = Inches(5.5), Inches(1.35)
    dashboards = [
        (im.get((4, 3)), "Dashboard Overview"),
        (im.get((4, 4)), "AI Business Line"),
    ]
    for i, (img, lbl) in enumerate(dashboards):
        dy = Inches(1.65) + Inches(i * 1.8)
        dash_frame(sl, img, dx, dy, dw, dh)
        tx(sl, lbl, dx, dy + dh + Inches(0.08), dw, Inches(0.25), sz=9, c=C_LT, al=PP_ALIGN.CENTER)

    ft(sl)


def s05_ai_line(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "The AI Business Line")
    tx(sl, "Our Secret Weapon", Inches(5.5), Inches(0.4), Inches(4), Inches(0.5), sz=20, c=C_YEL, b=True, al=PP_ALIGN.RIGHT)

    steps = [
        'Client taps "Refer a Friend" in their app',
        "Picks a contact from their phone",
        "Warm intro text sent from agent\u2019s AI number",
        "AI qualifies the lead via natural conversation",
        "AI books the appointment on agent\u2019s calendar",
        "Agent just shows up and closes",
    ]
    for i, s in enumerate(steps):
        sy = Inches(1.5) + Inches(i * 0.6)
        circ(sl, M, sy, Inches(0.4), C_ACC if i < 5 else C_YEL, str(i + 1), C_W, 14)
        tx(sl, s, M + Inches(0.55), sy, Inches(4.2), Inches(0.4), sz=14, c=C_OW, va=MSO_ANCHOR.MIDDLE)

    px = Inches(6.0)
    phone_frame(sl, im.get((5, 0)) or im.get((1, 1)), px, Inches(1.2), Inches(2.2), Inches(4.5))

    kx = Inches(8.7)
    lcrd(sl, kx, Inches(1.2), Inches(3.8), Inches(4.5), C_YEL)
    tx(sl, "KEY DIFFERENTIATOR", kx + Inches(0.3), Inches(1.4), Inches(3.2), Inches(0.35), sz=14, c=C_YEL, b=True)
    pts = [
        "The referral thinks they\u2019re texting the agent directly",
        "Natural, warm, personal \u2014 not a chatbot",
        "Auto-follows up on Day 2, 5 & 8",
        "Calls forward to agent\u2019s real phone",
    ]
    for j, p in enumerate(pts):
        bullet(sl, p, kx + Inches(0.3), Inches(2.0) + Inches(j * 0.55), Inches(3.2), C_YEL)

    ft(sl)


def s06_market(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "Market Opportunity")

    cw = Inches(3.6)
    ch = Inches(3.8)
    cy = Inches(1.5)

    markets = [
        ("TAM", "$4.5B", "All relationship-driven sales\nprofessionals across insurance,\nreal estate, financial advising,\nmortgage \u2014 1.2M+ insurance\nagents alone", C_ACC),
        ("SAM", "$530M", "900K+ licensed life & health\nagents in the US \u2014 independent\nAND captive \u2014 at $49/month\n($588/year)", C_BLUE),
        ("SOM", "$11M", "Symmetry Financial Group\u2019s\n6,000+ agents + adjacent\nIMO/FMO networks \u2014\nour built-in channel", C_YEL),
    ]

    for i, (label, amt, desc, ac) in enumerate(markets):
        cx = M + Inches(i * 3.9)
        acrd(sl, cx, cy, cw, ch, ac)
        tx(sl, label, cx, cy + Inches(0.25), cw, Inches(0.35), sz=14, c=ac, b=True, al=PP_ALIGN.CENTER)
        tx(sl, amt, cx, cy + Inches(0.65), cw, Inches(0.8), sz=48, c=C_W, b=True, al=PP_ALIGN.CENTER)
        ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx + Inches(0.4), cy + Inches(1.6), cw - Inches(0.8), Pt(1))
        ln.fill.solid(); ln.fill.fore_color.rgb = C_CARD_B; ln.line.fill.background()
        tx(sl, desc, cx + Inches(0.35), cy + Inches(1.8), cw - Inches(0.7), Inches(2.0), sz=12, c=C_LT)

    crd(sl, M, Inches(5.6), Inches(11.7), Inches(0.7), C_CARD, C_CARD_B)
    tx(sl, "\u2192  Beyond insurance: the white-label model scales to real estate agents, financial advisors, and mortgage brokers \u2014 millions of additional professionals.",
       M + Inches(0.2), Inches(5.65), Inches(11.3), Inches(0.6), sz=12, c=C_OW)

    ft(sl)


def s07_business(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "Business Model", "SaaS subscriptions with scarcity-driven launch tiers.")

    tiers = [
        ("Founding\nMembers", "FREE", "For Life", "50 spots", "7 CLAIMED", C_ACC),
        ("Charter\nMembers", "$25", "/month", "50 spots", "NEXT TIER", C_BLUE),
        ("Inner\nCircle", "$35", "/month", "50 spots", "UPCOMING", C_YEL),
        ("Standard", "$49", "/mo or $490/yr", "Unlimited", "TARGET", C_W),
    ]

    cw = Inches(2.65)
    cy = Inches(1.6)
    ch = Inches(3.6)

    for i, (name, price, period, spots, badge, ac) in enumerate(tiers):
        cx = M + Inches(i * 2.9)
        acrd(sl, cx, cy, cw, ch, ac)
        tx(sl, name, cx, cy + Inches(0.3), cw, Inches(0.6), sz=16, c=C_W, b=True, al=PP_ALIGN.CENTER)
        tx(sl, price, cx, cy + Inches(1.0), cw, Inches(0.6), sz=40, c=ac, b=True, al=PP_ALIGN.CENTER)
        tx(sl, period, cx, cy + Inches(1.65), cw, Inches(0.3), sz=13, c=C_OW, al=PP_ALIGN.CENTER)
        tx(sl, spots, cx, cy + Inches(2.1), cw, Inches(0.3), sz=12, c=C_LT, al=PP_ALIGN.CENTER)
        bw = Inches(1.6)
        bx = cx + (cw - bw) / 2
        badge_s = crd(sl, bx, cy + Inches(2.65), bw, Inches(0.35), ac if i == 0 else C_CARD, ac)
        tx(sl, badge, bx, cy + Inches(2.65), bw, Inches(0.35), sz=10, c=C_W if i == 0 else ac, b=True, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    crd(sl, M, Inches(5.55), Inches(11.7), Inches(0.65), C_CARD, C_ACC)
    tx(sl, "\u2705  1 saved policy or 1 referral = $1,200+ annual value.  At $49/mo, that\u2019s instant 2x+ ROI for the agent.",
       M + Inches(0.2), Inches(5.6), Inches(11.3), Inches(0.55), sz=14, c=C_W, b=True, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    ft(sl)


def s08_traction(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "Traction & Milestones")

    stats = [
        ("LIVE", "iOS & Google Play", C_ACC),
        ("7", "Founding Members", C_BLUE),
        ("1", "Solo Founder\nBuilt Everything", C_YEL),
        ("$0", "External Funding\nTo Date", C_CORAL),
    ]

    for i, (v, lbl, ac) in enumerate(stats):
        cx = M + Inches(i * 2.9)
        acrd(sl, cx, Inches(1.4), Inches(2.65), Inches(1.4), ac)
        tx(sl, v, cx, Inches(1.6), Inches(2.65), Inches(0.55), sz=36, c=ac, b=True, al=PP_ALIGN.CENTER)
        tx(sl, lbl, cx, Inches(2.15), Inches(2.65), Inches(0.5), sz=11, c=C_OW, al=PP_ALIGN.CENTER)

    tx(sl, "What We\u2019ve Built", M, Inches(3.1), Inches(11), Inches(0.35), sz=20, c=C_W, b=True)
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(3.45), Inches(11.7), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_CARD_B; ln.line.fill.background()

    items = [
        ("Full platform: web dashboard, iOS app, Android app, AI business line, backend infrastructure", C_ACC),
        ("Live on App Store & Google Play with 7 founding member agents signed up", C_ACC),
        ("DBA registered, IP assigned to Brainstorm Labs LLC, entity formation complete", C_ACC),
        ("AI conservation alerts, referral automation, 7+ annual touchpoints all functional", C_ACC),
    ]
    for j, (t, ac) in enumerate(items):
        by = Inches(3.65) + Inches(j * 0.45)
        circ(sl, M + Inches(0.1), by + Inches(0.04), Inches(0.18), ac, "\u2713", C_W, 10)
        tx(sl, t, M + Inches(0.45), by, Inches(11), Inches(0.35), sz=13, c=C_OW)

    crd(sl, M, Inches(5.55), Inches(11.7), Inches(0.65), C_PRI_DK, C_ACC)
    tx(sl, "\u2192  NEXT: Fill 50 founding seats  \u2192  Launch paid tiers  \u2192  Reach $5K+ MRR by Q3 2026",
       M + Inches(0.2), Inches(5.6), Inches(11.3), Inches(0.55), sz=14, c=C_YEL, b=True, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    ft(sl)


def s09_gtm(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "Go-to-Market Strategy")

    phases = [
        ("NOW", "Founder-Led", [
            "Direct outreach within Symmetry Financial Group (6,000+ agents)",
            "Insurance agent Facebook & WhatsApp communities",
            "Free founding seats build case studies & social proof",
        ], C_ACC),
        ("2026\u201327", "Channel Partners", [
            "IMO/FMO leaders distribute to their downlines",
            'SEO: "insurance chargebacks", "agent referral system"',
            "Conference presence & community content",
        ], C_BLUE),
        ("2027+", "Market Expansion", [
            "Real estate, financial advisors, mortgage brokers",
            "AI business line model requires minimal adaptation",
            "Enterprise platform licensing for large organizations",
        ], C_YEL),
    ]

    cw = Inches(3.6)
    ch = Inches(4.5)
    cy = Inches(1.5)

    for i, (phase, title, pts, ac) in enumerate(phases):
        cx = M + Inches(i * 3.9)
        acrd(sl, cx, cy, cw, ch, ac)
        badge_w = Inches(1.1)
        badge_s = crd(sl, cx + Inches(0.25), cy + Inches(0.25), badge_w, Inches(0.3), ac, ac)
        tx(sl, phase, cx + Inches(0.25), cy + Inches(0.25), badge_w, Inches(0.3), sz=10, c=C_W, b=True, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
        tx(sl, title, cx + Inches(0.25), cy + Inches(0.7), cw - Inches(0.5), Inches(0.35), sz=20, c=C_W, b=True)
        for j, p in enumerate(pts):
            bullet(sl, p, cx + Inches(0.25), cy + Inches(1.35) + Inches(j * 0.7), cw - Inches(0.5), ac)

    ft(sl)


def s10_competition(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "Competitive Landscape", "No competitor combines all three revenue streams with AI in one platform.")

    features = [
        "Branded client app",
        "AI business line",
        "Automated retention",
        "One-tap referrals",
        "Conservation alerts",
        "Anniversary rewrites",
    ]
    companies = [
        ("AgentForLife", [True]*6, "$49"),
        ("AgencyBloc",   [False, False, "P", False, False, False], "$65+"),
        ("Radiusbob",    [False, False, "P", False, False, False], "$34+"),
        ("Salesforce",   [False]*6, "$300+"),
    ]

    gx = M
    gy = Inches(1.55)
    fw = Inches(2.4)
    cw_col = Inches(2.1)
    rh = Inches(0.38)
    hdr_h = Inches(0.42)

    for ci, (name, feats, price) in enumerate(companies):
        cx = gx + fw + Inches(ci * 2.2)
        h_fill = C_PRI_DK if ci == 0 else C_CARD
        h_brd = C_ACC if ci == 0 else C_CARD_B
        hdr = crd(sl, cx, gy, cw_col, hdr_h, h_fill, h_brd)
        color = C_ACC if ci == 0 else C_W
        tx(sl, name, cx, gy, cw_col, hdr_h, sz=11, c=color, b=True, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    for fi, feat in enumerate(features):
        ry = gy + hdr_h + Inches(0.05) + Inches(fi * 0.42)
        row_bg = C_CARD if fi % 2 == 0 else RGBColor(0x08, 0x28, 0x25)
        crd(sl, gx, ry, fw, rh, row_bg, C_CARD_B)
        tx(sl, feat, gx + Inches(0.15), ry, fw - Inches(0.2), rh, sz=11, c=C_OW, va=MSO_ANCHOR.MIDDLE)

        for ci, (_, feats, _) in enumerate(companies):
            cx = gx + fw + Inches(ci * 2.2)
            cell_bg = C_CARD if fi % 2 == 0 else RGBColor(0x08, 0x28, 0x25)
            crd(sl, cx, ry, cw_col, rh, cell_bg, C_CARD_B)
            val = feats[fi]
            if val is True:
                tx(sl, "\u2713", cx, ry, cw_col, rh, sz=16, c=C_ACC, b=True, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
            elif val == "P":
                tx(sl, "Partial", cx, ry, cw_col, rh, sz=10, c=C_DIM, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
            else:
                tx(sl, "\u2717", cx, ry, cw_col, rh, sz=14, c=C_CORAL, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    price_y = gy + hdr_h + Inches(0.05) + Inches(6 * 0.42)
    crd(sl, gx, price_y, fw, rh, C_PRI_DK, C_CARD_B)
    tx(sl, "Price / month", gx + Inches(0.15), price_y, fw - Inches(0.2), rh, sz=11, c=C_W, b=True, va=MSO_ANCHOR.MIDDLE)
    for ci, (_, _, price) in enumerate(companies):
        cx = gx + fw + Inches(ci * 2.2)
        c_fill = C_PRI_DK if ci == 0 else C_CARD
        crd(sl, cx, price_y, cw_col, rh, c_fill, C_CARD_B)
        pc = C_ACC if ci == 0 else C_W
        tx(sl, price, cx, price_y, cw_col, rh, sz=12, c=pc, b=True, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    chart_y = Inches(4.7)
    chart_h = Inches(1.7)
    chart_data = CategoryChartData()
    chart_data.categories = ["AgentForLife", "AgencyBloc", "Radiusbob", "Salesforce"]
    chart_data.add_series("Features", (6, 1, 1, 0))

    cf = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, M, chart_y, Inches(11.7), chart_h, chart_data)
    chart = cf.chart
    chart.has_legend = False
    chart.has_title = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(12); dl.font.bold = True; dl.font.name = FONT
    dl.font.color.rgb = C_W
    dl.number_format = '0"/6"'

    series = chart.series[0]
    for i, color in enumerate([C_ACC, C_DIM, C_DIM, C_DIM]):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = color

    plot.gap_width = 100

    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(10)
    ca.tick_labels.font.color.rgb = C_LT
    ca.tick_labels.font.name = FONT
    ca.format.line.fill.background()
    ca.has_major_gridlines = False

    va = chart.value_axis
    va.visible = False
    va.has_major_gridlines = False
    va.format.line.fill.background()

    try:
        cs = chart._chartSpace
        for spPr_tag in (qn("c:spPr"),):
            for el in cs.findall(spPr_tag):
                cs.remove(el)
        spPr = cs.makeelement(qn("c:spPr"), {})
        spPr.append(spPr.makeelement(qn("a:noFill"), {}))
        cs.append(spPr)
        c_chart = cs.find(qn("c:chart"))
        pa = c_chart.find(qn("c:plotArea"))
        for el in pa.findall(qn("c:spPr")):
            pa.remove(el)
        pa_spPr = pa.makeelement(qn("c:spPr"), {})
        pa_spPr.append(pa_spPr.makeelement(qn("a:noFill"), {}))
        pa.append(pa_spPr)
    except:
        pass

    ft(sl)


def s11_math(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tx(sl, "The Math That Sells Itself", M, Inches(0.3), Inches(11), Inches(0.6), sz=36, c=C_W, b=True)

    cw = Inches(5.5)
    cy = Inches(1.15)
    ch = Inches(4.8)

    crd(sl, M, cy, cw, ch, C_RED_CARD, C_RED_BRD)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, cy, cw, Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_CORAL; bar.line.fill.background()
    tx(sl, "WITHOUT AgentForLife", M, cy + Inches(0.15), cw, Inches(0.35), sz=16, c=C_CORAL, b=True, al=PP_ALIGN.CENTER)
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, M + Inches(0.4), cy + Inches(0.55), cw - Inches(0.8), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_RED_BRD; ln.line.fill.background()

    without = [
        ("$4,000/mo", "spent on cold leads"),
        ("160", "strangers who may not answer"),
        ("10\u201330%", "contact rate on cold leads"),
        ("15\u201320%", "of closed policies lapse \u2192 chargebacks"),
        ("$0", "revenue from existing clients"),
    ]
    for j, (v, d) in enumerate(without):
        ry = cy + Inches(0.75) + Inches(j * 0.75)
        tx(sl, v, M + Inches(0.3), ry, Inches(1.8), Inches(0.4), sz=18, c=C_CORAL, b=True)
        tx(sl, d, M + Inches(2.2), ry, Inches(2.8), Inches(0.4), sz=12, c=C_OW, va=MSO_ANCHOR.MIDDLE)

    rx = Inches(7.0)
    crd(sl, rx, cy, cw, ch, C_CARD, C_ACC)
    bar2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, cy, cw, Pt(4))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = C_ACC; bar2.line.fill.background()
    tx(sl, "WITH AgentForLife", rx, cy + Inches(0.15), cw, Inches(0.35), sz=16, c=C_ACC, b=True, al=PP_ALIGN.CENTER)
    ln2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx + Inches(0.4), cy + Inches(0.55), cw - Inches(0.8), Pt(1))
    ln2.fill.solid(); ln2.fill.fore_color.rgb = C_CARD_B; ln2.line.fill.background()

    with_afl = [
        ("$49/mo", "less than ONE cold lead"),
        ("7+", "touchpoints keep clients retained"),
        ("50%+", "close rate on warm referrals"),
        ("$1,200+", "saved per prevented chargeback"),
        ("\u221E", "referrals from your existing book"),
    ]
    for j, (v, d) in enumerate(with_afl):
        ry = cy + Inches(0.75) + Inches(j * 0.75)
        tx(sl, v, rx + Inches(0.3), ry, Inches(1.8), Inches(0.4), sz=18, c=C_ACC, b=True)
        tx(sl, d, rx + Inches(2.2), ry, Inches(2.8), Inches(0.4), sz=12, c=C_OW, va=MSO_ANCHOR.MIDDLE)

    crd(sl, M, Inches(6.1), Inches(11.7), Inches(0.5), C_PRI_DK, C_ACC)
    tx(sl, "For the cost of ONE lead, triple your income from clients you already have.",
       M, Inches(6.12), Inches(11.7), Inches(0.5), sz=15, c=C_W, b=True, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    ft(sl)


def s12_moat(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "Why Us \u2014 Our Moat", "Why another company can\u2019t just copy this.")

    moats = [
        (C_ACC, "01", "Founder IS the Customer",
         "I\u2019m an active insurance agent writing policies every week. Every feature is battle-tested against my own book of business before it ships. A tech company would need years of discovery to learn what I already know."),
        (C_BLUE, "02", "Built-In Distribution",
         "Already inside Symmetry Financial Group (6,000+ agents). The IMO/FMO model means one partnership = thousands of agents. A competitor starts from zero relationships in a trust-based industry."),
        (C_YEL, "03", "Compounding Switching Costs",
         "Once clients download the branded app and the AI number is on business cards, the agent can\u2019t leave without disrupting every client relationship. Every month makes it stickier."),
        (C_CORAL, "04", "Conversation Intelligence",
         "Every referral conversation reveals what language, timing, and follow-up cadence converts. This playbook compounds over time and can\u2019t be replicated without thousands of real conversations."),
    ]

    cw = Inches(5.6)
    ch = Inches(2.05)

    for i, (ac, icon, title, desc) in enumerate(moats):
        cx = M if i % 2 == 0 else Inches(6.9)
        cy = Inches(1.55) + Inches((i // 2) * 2.3)
        lcrd(sl, cx, cy, cw, ch, ac)
        circ(sl, cx + Inches(0.25), cy + Inches(0.2), Inches(0.45), ac, icon, C_W, 18)
        tx(sl, title, cx + Inches(0.85), cy + Inches(0.25), cw - Inches(1.1), Inches(0.35), sz=17, c=C_W, b=True)
        tx(sl, desc, cx + Inches(0.25), cy + Inches(0.75), cw - Inches(0.5), Inches(1.2), sz=11, c=C_LT)

    ft(sl)


def s13_team(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "The Team")

    lcrd(sl, M, Inches(1.3), Inches(5.5), Inches(5.0), C_ACC)
    tx(sl, "Daniel Roberts", M + Inches(0.3), Inches(1.5), Inches(5), Inches(0.4), sz=24, c=C_W, b=True)
    tx(sl, "Founder & CEO", M + Inches(0.3), Inches(1.95), Inches(5), Inches(0.3), sz=14, c=C_ACC)

    bio = [
        "Active insurance agent \u2014 built from personal pain",
        "Self-taught developer \u2014 solo-built entire platform",
        "JD, Saint Louis University School of Law",
        "BS, Miami University",
        "Former brand strategist at Grain (STL)",
        "Born & raised in St. Louis, returned Sept 2025",
    ]
    for j, b in enumerate(bio):
        bullet(sl, b, M + Inches(0.3), Inches(2.55) + Inches(j * 0.42), Inches(4.8), C_ACC)

    rx = Inches(7.0)
    crd(sl, rx, Inches(1.3), Inches(5.5), Inches(2.0), C_CARD, C_CARD_B)
    tx(sl, "Why St. Louis", rx + Inches(0.2), Inches(1.4), Inches(5), Inches(0.35), sz=18, c=C_YEL, b=True)
    tx(sl, "3rd-generation STL entrepreneur. Chose to return\nafter 10 years in NYC & LA. Lives & works in Cortex\nInnovation District. Plans to hire 2\u20133 local team\nmembers by 2027.",
       rx + Inches(0.2), Inches(1.85), Inches(5), Inches(1.3), sz=12, c=C_OW)

    crd(sl, rx, Inches(3.55), Inches(5.5), Inches(2.75), C_CARD, C_CARD_B)
    tx(sl, "Mentor Network", rx + Inches(0.2), Inches(3.65), Inches(5), Inches(0.35), sz=16, c=C_ACC, b=True)

    mentors = [
        "Maxine Clark \u2014 Build-A-Bear (Founder)",
        "Jeffrey Winters \u2014 Abstrakt Marketing",
        "Neal Fenster \u2014 Enterprise Medical Recruiting",
        "Ben Weiss \u2014 Daytona Street Capital",
        "Dan & Sarah Mirth \u2014 Artifox (2025 ARCH Grantee)",
    ]
    for j, m in enumerate(mentors):
        bullet(sl, m, rx + Inches(0.25), Inches(4.1) + Inches(j * 0.38), Inches(5), C_YEL, C_OW, 11)

    ft(sl)


def s14_impact(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    ttl(sl, "St. Louis Impact Plan", "Building AgentForLife in St. Louis, for St. Louis \u2014 and beyond.")

    years = [
        ("Year 1", "Foundation", [
            "Headquarters in Cortex Innovation District",
            "Fill founding tiers, validate product-market fit",
            "Engage STL community (1MC, ITEN, T-REX)",
            "Contribute to STL fintech/insurtech cluster",
        ], C_ACC),
        ("Year 2", "Growth & Hiring", [
            "Hire 2\u20133 local team members",
            "Reach 500+ paying agent subscribers",
            "Partner with STL companies for co-marketing",
            "Generate $200K+ ARR from STL HQ",
        ], C_BLUE),
        ("Year 3+", "Scale & Leadership", [
            "Expand team to 5\u20138 STL employees",
            "Become leading insurtech for agent retention",
            "Mentor future ARCH Grants applicants",
            "Establish as a St. Louis success story",
        ], C_YEL),
    ]

    cw = Inches(3.6)
    ch = Inches(4.5)
    cy = Inches(1.5)

    for i, (yr, title, pts, ac) in enumerate(years):
        cx = M + Inches(i * 3.9)
        acrd(sl, cx, cy, cw, ch, ac)
        badge_w = Inches(1.0)
        b_s = crd(sl, cx + Inches(0.2), cy + Inches(0.2), badge_w, Inches(0.3), ac, ac)
        tx(sl, yr, cx + Inches(0.2), cy + Inches(0.2), badge_w, Inches(0.3), sz=10, c=C_W, b=True, al=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
        tx(sl, title, cx + Inches(0.2), cy + Inches(0.65), cw - Inches(0.4), Inches(0.35), sz=20, c=C_W, b=True)
        for j, p in enumerate(pts):
            bullet(sl, p, cx + Inches(0.2), cy + Inches(1.25) + Inches(j * 0.65), cw - Inches(0.4), ac)

    ft(sl)


def s15_ask(prs, im):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)

    tx(sl, "The Ask", M, Inches(0.3), Inches(11), Inches(0.5), sz=36, c=C_W, b=True, al=PP_ALIGN.CENTER)
    tx(sl, "$75,000", Inches(0.5), Inches(0.9), Inches(12.3), Inches(1.0), sz=72, c=C_ACC, b=True, al=PP_ALIGN.CENTER)
    tx(sl, "Equity-Free ARCH Grant", Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.45), sz=22, c=C_OW, al=PP_ALIGN.CENTER)

    chart_data = CategoryChartData()
    chart_data.categories = [
        "Founder Salary\n& Runway (35%)",
        "Customer\nAcquisition (30%)",
        "AI &\nInfrastructure (25%)",
        "Legal &\nOperations (10%)",
    ]
    chart_data.add_series("Allocation", (35, 30, 25, 10))

    cf = sl.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(0.8), Inches(2.4), Inches(5.5), Inches(3.8), chart_data
    )
    chart = cf.chart
    chart.has_legend = True
    chart.has_title = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(11); dl.font.bold = True; dl.font.name = FONT
    dl.font.color.rgb = C_W
    dl.number_format = '0"%"'

    series = chart.series[0]
    for i, color in enumerate([C_ACC, C_BLUE, C_YEL, C_CORAL]):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = color

    legend = chart.legend
    legend.include_in_layout = False
    legend.position = XL_LEGEND_POSITION.BOTTOM
    legend.font.size = Pt(10); legend.font.name = FONT
    legend.font.color.rgb = C_LT

    try:
        cs = chart._chartSpace
        for el in cs.findall(qn("c:spPr")):
            cs.remove(el)
        spPr = cs.makeelement(qn("c:spPr"), {})
        spPr.append(spPr.makeelement(qn("a:noFill"), {}))
        cs.append(spPr)
        c_chart = cs.find(qn("c:chart"))
        pa = c_chart.find(qn("c:plotArea"))
        for el in pa.findall(qn("c:spPr")):
            pa.remove(el)
        pa_sp = pa.makeelement(qn("c:spPr"), {})
        pa_sp.append(pa_sp.makeelement(qn("a:noFill"), {}))
        pa.append(pa_sp)
    except:
        pass

    alloc = [
        ("35%", "Founder salary & runway", "Dedicated focus on growth, not side income", C_ACC),
        ("30%", "Customer acquisition", "Marketing, content, IMO/FMO outreach", C_BLUE),
        ("25%", "AI & infrastructure", "Scale AI business line, hosting, API costs", C_YEL),
        ("10%", "Legal & operations", "Trademark, compliance, business ops", C_CORAL),
    ]

    rx = Inches(7.0)
    for i, (pct, name, desc, ac) in enumerate(alloc):
        ry = Inches(2.6) + Inches(i * 0.9)
        crd(sl, rx, ry, Inches(5.5), Inches(0.75), C_CARD, C_CARD_B)
        bar_w = Inches(float(pct.replace("%", "")) / 100 * 5.5)
        bar_s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, ry, bar_w, Inches(0.75))
        bar_s.fill.solid(); bar_s.fill.fore_color.rgb = ac
        bar_s.line.fill.background()
        try:
            spPr = bar_s._element.find(".//" + qn("a:solidFill"))
            if spPr is not None:
                srgb = spPr.find(qn("a:srgbClr"))
                if srgb is not None:
                    etree.SubElement(srgb, qn("a:alpha"), {"val": "25000"})
        except:
            pass
        tx(sl, pct, rx + Inches(0.15), ry, Inches(0.7), Inches(0.75), sz=18, c=ac, b=True, va=MSO_ANCHOR.MIDDLE)
        tx(sl, name, rx + Inches(0.85), ry + Inches(0.05), Inches(2.5), Inches(0.35), sz=13, c=C_W, b=True)
        tx(sl, desc, rx + Inches(0.85), ry + Inches(0.38), Inches(4), Inches(0.3), sz=10, c=C_LT)

    ft(sl, "agentforlife.app  |  daniel@agentforlife.app  |  (314) 363-4922  |  St. Louis, MO")


# ═══════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    src = "/Users/danielroberts/Downloads/AgentForLife_ARCH_Grants_Pitch_Deck_FINAL.pptx"
    out = "/Users/danielroberts/Downloads/AgentForLife_ARCH_Grants_Premium.pptx"

    tmpdir = tempfile.mkdtemp(prefix="afl_deck_")
    print(f"Extracting images to {tmpdir}...")
    im = extract_images(src, tmpdir)
    print(f"  Found {len(im)} images")

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    builders = [
        s01_title, s02_problem, s03_solution, s04_product, s05_ai_line,
        s06_market, s07_business, s08_traction, s09_gtm, s10_competition,
        s11_math, s12_moat, s13_team, s14_impact, s15_ask,
    ]

    for i, builder in enumerate(builders):
        print(f"  Building slide {i+1}/15...")
        builder(prs, im)

    prs.save(out)
    print(f"\nDone! Saved to:\n  {out}")

if __name__ == "__main__":
    main()
